import io
import logging
import os
from distutils.util import strtobool

import langchain
import openpyxl
import streamlit as st
import transformers
from chromadb.config import Settings
from dotenv import load_dotenv
from langchain import hub
from langchain.agents import AgentExecutor
from langchain.agents import AgentType
from langchain.agents import BaseSingleActionAgent
from langchain.agents import LLMSingleActionAgent
from langchain.agents import Tool
from langchain.agents import create_openai_tools_agent
from langchain.agents import create_react_agent
from langchain.agents import initialize_agent
from langchain.callbacks.manager import AsyncCallbackManagerForToolRun
from langchain.callbacks.manager import CallbackManagerForToolRun
from langchain.callbacks.tracers import ConsoleCallbackHandler
from langchain.chains import ConversationChain
from langchain.chains import LLMChain
from langchain.chains import RetrievalQAWithSourcesChain
from langchain.chains.conversation.prompt import PROMPT
from langchain.chains.retrieval_qa.base import RetrievalQA
from langchain.memory import ConversationBufferMemory
from langchain.memory import ConversationBufferWindowMemory
from langchain.prompts import load_prompt
from langchain.prompts.chat import AIMessagePromptTemplate
from langchain.prompts.chat import ChatPromptTemplate
from langchain.prompts.chat import HumanMessagePromptTemplate
from langchain.prompts.chat import MessagesPlaceholder
from langchain.prompts.chat import SystemMessagePromptTemplate
from langchain.pydantic_v1 import BaseModel
from langchain.pydantic_v1 import Field
from langchain.retrievers.self_query.base import SelfQueryRetriever
from langchain.retrievers.web_research import WebResearchRetriever
from langchain.schema import AIMessage
from langchain.schema import HumanMessage
from langchain.schema import SystemMessage
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.tools import BaseTool
from langchain.tools import StructuredTool
from langchain.tools import Tool
from langchain_community.callbacks import get_openai_callback
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain_community.chat_models import ChatOllama
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import Ollama
from langchain_community.utilities import GoogleSearchAPIWrapper
from langchain_community.vectorstores import Chroma
from langchain_core.chat_history import BaseChatMessageHistory
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_openai import ChatOpenAI
from langchain_openai.llms.base import OpenAI
from pptx import Presentation
from PyPDF2 import PdfReader

load_dotenv()
langchain.debug = True

api_base = os.environ['API_BASE']
vectordb_host = os.environ['VECTORDB_HOST']
vectordb_port = os.environ['VECTORDB_PORT']
global_cache_dir = 'models'
global_log_verbosity = bool(strtobool(os.environ['LOG_VERBOSE']))
logger = logging.getLogger(st.__name__)
logger.setLevel(logging.INFO)
model_id_gemma = 0
model_id_gpt3 = 1
language_ja = 0
language_en = 1

@st.cache_resource
def get_local_embeddings():
    """
    Get a local embedding model.

    Returns
    -------
    embeddings : model instance
    """
    model_kwargs = {'device': 'cpu'}
    return HuggingFaceEmbeddings(
        model_name = 'intfloat/multilingual-e5-large',
        model_kwargs = model_kwargs,
        cache_folder=global_cache_dir)

@st.cache_resource
def get_vectordb_web():
    """
    Get persistent vector db for web scrapping.

    Returns
    -------
    persistent vector db for web search : db instance
    """
    model_kwargs = {'device': 'cpu'}
    embeddings = get_local_embeddings()
    return Chroma(collection_name='chroma_web', embedding_function=embeddings)# persist_directory='chroma_web')

@st.cache_resource
def get_vectordb_doc():
    """
    Get persistent vector db for documents
    Returns
    -------
    persistent vector db for documents : db instance
    """
    model_kwargs = {'device': 'cpu'}
    embeddings = get_local_embeddings()
    return Chroma(collection_name='chroma_doc', embedding_function=embeddings)#, persist_directory='chroma_doc')

@st.cache_resource
def get_vectordb_server(collection_name):
    """
    Get a client server model db.
    Returns
    -------
    persistent vector : db instance
    """
    embeddings = get_local_embeddings()
    client_settings = Settings(
        chroma_api_impl="chromadb.api.fastapi.FastAPI",
        chroma_server_host=vectordb_host,
        chroma_server_http_port=vectordb_port
    )
    return Chroma(
        collection_name=collection_name,
        client_settings=client_settings,
        embedding_function=embeddings
    )

@st.cache_resource
def get_google_api():
    return GoogleSearchAPIWrapper()

def get_session_history(session_id: str) -> BaseChatMessageHistory:
    return st.session_state.history

#@st.cache_resource
#def get_trans_ja_en():
#    model = MarianMTModel.from_pretrained(pretrained_model_name_or_path='staka/fugumt-ja-en', cache_dir=global_cache_dir)
#    tokenizer = AutoTokenizer.from_pretrained(pretrained_model_name_or_path='staka/fugumt-ja-en', cache_dir=global_cache_dir)
#    return pipeline('translation', model=model, tokenizer=tokenizer)

#@st.cache_resource
#def get_trans_en_ja():
#    model = MarianMTModel.from_pretrained(pretrained_model_name_or_path='staka/fugumt-en-ja', cache_dir=global_cache_dir)
#    tokenizer = AutoTokenizer.from_pretrained(pretrained_model_name_or_path='staka/fugumt-en-ja', cache_dir=global_cache_dir)
#    return pipeline('translation', model=model, tokenizer=tokenizer)

#@st.cache_resource
#def get_fasttext():
#    return fasttext.load_model('models/lid.176.bin')

### Prompt templates
####################
CHAT_TEMPLATE_JA = "あなたは誠実で優秀な日本人のAIです。ユーザのメッセージに対して、短く豊かな返信文を作成してください。"
CHAT_PROMPT_JA = ChatPromptTemplate(
    messages=[
        SystemMessagePromptTemplate.from_template(
            CHAT_TEMPLATE_JA
        ),
        # The `variable_name` here is what must align with memory
        MessagesPlaceholder(variable_name="history"),
        HumanMessagePromptTemplate.from_template("{input}")
    ]
)

CHAT_TEMPLATE_EN = "You are a honest and excelent AI. Please compose short and rich replies to human messages."
CHAT_PROMPT_EN = ChatPromptTemplate(
    messages=[
        SystemMessagePromptTemplate.from_template(
            CHAT_TEMPLATE_EN
        ),
        # The `variable_name` here is what must align with memory
        MessagesPlaceholder(variable_name="history"),
        HumanMessagePromptTemplate.from_template("{input}")
    ]
)

def get_chat_prompt():
    if model_id_gemma == st.session_state.model_id:
        if language_ja == st.session_state.language:
            return CHAT_PROMPT_JA
        else:
            return CHAT_PROMPT_EN

    if model_id_gpt3 == st.session_state.model_id:
        return PROMPT

SEARCH_QUERY_TEMPLATE_JA =\
"""あなたはGoogle検索の結果を改善するアシスタントです。
次の質問に似たGoogle検索のクエリを3つ作ってください。出力クエリは番号付きのリストで返してください。: {question}"""

SEARCH_QUERY_PROMPT_JA = PromptTemplate(
    input_variables=["question"],
    template=SEARCH_QUERY_TEMPLATE_JA,
)

SEARCH_QUERY_TEMPLATE_EN =\
"""You are an assistant tasked with improving Google search results.
Generate THREE Google search queries that are similar to this question. The output should be a numbered list of queries.: {question}"""

SEARCH_QUERY_PROMPT_EN = PromptTemplate(
    input_variables=["question"],
    template=SEARCH_QUERY_TEMPLATE_EN,
)

def get_search_query_prompt():
    if model_id_gpt3 == st.session_state.model_id:
        return None

    if language_ja == st.session_state.language:
        return SEARCH_QUERY_PROMPT_JA
    else:
        return SEARCH_QUERY_PROMPT_EN

QA_BASE_PROMPT_JA = PromptTemplate(input_variables=['page_content', 'source'], template='Content: {page_content}\nSource: {source}')
QA_BASE_PROMPT_EN = PromptTemplate(input_variables=['page_content', 'source'], template='Content: {page_content}\nSource: {source}')

QA_QUESTION_PROMPT_JA = PromptTemplate(
    input_variables=['context', 'question'],
    template='長い文章の次の部分を使って、質問に答えるために関連する文章があるか確認してください。\n関連するテキストをそのまま返してください。\n{context}\n質問: {question}\nもしあれば、関連する文章:')

QA_QUESTION_PROMPT_EN = PromptTemplate(
    input_variables=['context', 'question'],
    template='Use the following portion of a long document to see if any of the text is relevant to answer the question. \nReturn any relevant text verbatim.\n{context}\nQuestion: {question}\nRelevant text, if any:')

QA_COMBINE_TEMPLATE_JA =\
"""Given the following extracted parts of a long document and a question, create a final answer with references ("SOURCES").
If you don\'t know the answer, just say that you don\'t know. Don\'t try to make up an answer.
ALWAYS return a "SOURCES" part in your answer.
SOURCES:
{summaries}

QUESTION: {question}

FINAL ANSWER:
"""
QA_COMBINE_PROMPT_JA = PromptTemplate(input_variables=['question', 'summaries'], template=QA_COMBINE_TEMPLATE_JA)

QA_COMBINE_TEMPLATE_EN =\
"""Given the following extracted parts of a long document and a question, create a final answer with references ("SOURCES").
If you don\'t know the answer, just say that you don\'t know. Don\'t try to make up an answer.
ALWAYS return a "SOURCES" part in your answer.
SOURCES:
{summaries}

QUESTION: {question}

FINAL ANSWER:
"""
QA_COMBINE_PROMPT_EN = PromptTemplate(input_variables=['question', 'summaries'], template=QA_COMBINE_TEMPLATE_EN)

def get_qa_source_prompt():
    if model_id_gpt3 == st.session_state.model_id:
        return None

    if language_ja == st.session_state.language:
        return (QA_BASE_PROMPT_JA, QA_QUESTION_PROMPT_JA, QA_COMBINE_PROMPT_JA) 
    else:
        return (QA_BASE_PROMPT_EN, QA_QUESTION_PROMPT_EN, QA_COMBINE_PROMPT_EN)

TEXT_QA_PROMPT_JA =\
"""次の文脈情報があります。

{context}

事前情報を使わず、与えられた文脈情報を使って、次の質問に答えてください。
質問: {question}
"""

TEXT_QA_PROMPT_JA = PromptTemplate(
    input_variables=['context', 'question'],
    template = TEXT_QA_PROMPT_JA)

TEXT_QA_PROMPT_EN =\
"""Context information is below.

{context}

Given the context information and not prior knowledge,
answer the question: {question}
"""

TEXT_QA_PROMPT_EN = PromptTemplate(
    input_variables=['context', 'question'],
    template = TEXT_QA_PROMPT_EN)

def get_qa_prompt():
    if model_id_gpt3 == st.session_state.model_id:
        return None

    if language_ja == st.session_state.language:
        return TEXT_QA_PROMPT_JA 
    else:
        return TEXT_QA_PROMPT_EN

def split_text(text:str):
    splitter = RecursiveCharacterTextSplitter(chunk_size=400, chunk_overlap=40)
    return splitter.split_text(text)

def load_file(file:io.BytesIO):
    content = None
    if 'application/pdf' == file.type:
        pdf_reader = PdfReader(file)
        content = '\n'.join([page.extract_text() for page in pdf_reader.pages])
    elif 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' == file.type:
        book = openpyxl.load_workbook(file)
        sheets = book.sheetnames
        content = []
        for sheet_name in sheets:
            sheet = book[sheet_name]
            cell_list = [f"{cell.value}" for cells in tuple(sheet.columns) for cell in cells]
            content.extend(cell_list)
    elif 'application/vnd.openxmlformats-officedocument.presentationml.presentation' == file.type:
        presentation = Presentation(file)
        content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for par in shape.text_frame.paragraphs:
                    for run in par.runs:
                        content.append(run.text)
    elif 'text/plain' == file.type or 'text/csv' == file.type:
        content = file.read().decode('UTF-8')
    else:
        return None
    file.seek(0)
    return content

def embeddig_file(file:io.BytesIO):
    with st.spinner('Embedding ...'):
        text = load_file(file)
        chunks = split_text(text)
        vector_db = get_vectordb_doc()
        vector_db.add_texts(chunks)

def clear_session_state():
    if 'history' in st.session_state:
        st.session_state.history.clear()
    else:
        st.session_state.history = ChatMessageHistory()
    if 'messages' in st.session_state:
        st.session_state.messages = []
    else:
        st.session_state.messages = []
    st.session_state.model = model_id_gemma
    st.session_state.language = language_ja 
    st.session_state.upload_file = None

def init_page():
    st.set_page_config(
            page_title='Chat Agent',
            page_icon='😊',
            layout='wide',
    )
    st.header('😊 Chat Agent 😊')
    st.sidebar.title("Options")

def choose_model(chat_mode=False):
    if model_id_gemma == st.session_state.model_id:
        model_name = 'eramax/gemma-7b-it:q4_k_m'
        if chat_mode:
            st.session_state.model = ChatOllama(
                model=model_name,
                base_url=api_base,
                temperature=st.session_state.temperature,
                num_ctx=8192,
                verbose=global_log_verbosity)
        else:
            st.session_state.model = Ollama(
                model=model_name,
                base_url=api_base,
                temperature=st.session_state.temperature,
                num_ctx=8192,
                verbose=global_log_verbosity)
    else:
        model_name = 'gpt-3.5-turbo'
        st.session_state.model = OpenAI(
            temperature=st.session_state.temperature,
            model_name=model_name,
            max_tokens=4096,
            verbose=global_log_verbosity)

def init_options():
    ### Choose a model
    container = st.sidebar.container(border=True)
    model = container.radio('モデル:', ('gemma-7b-it', 'GPT-3.5'))
    st.session_state.temperature = container.slider('Temerature:', min_value=0.0, max_value=1.0, value=0.1)

    if 'gemma-7b-it' == model:
        st.session_state.model_id = model_id_gemma
    else:
        st.session_state.model_id = model_id_gpt3

    ### Choose a language
    container = st.sidebar.container(border=True)
    languages = {0: "日本語", 1: "English"}
    def language_format_func(option):
        return languages[option]

    st.session_state.language = container.selectbox("Language:", options=list(languages.keys()), format_func=language_format_func)



def init_messages():
    clear_button = st.sidebar.button('Clear Conversation', key='clear')
    if clear_button or 'history' not in st.session_state:
        clear_session_state()

def invoke_chat_chain(user_input:str):
    choose_model(chat_mode=True)
    prompt = get_chat_prompt()
    chain = prompt | st.session_state.model | StrOutputParser()
    with_message_history = RunnableWithMessageHistory(
        chain,
        get_session_history,
        input_messages_key='input',
        history_messages_key='history',
        )
    return with_message_history.invoke({'input': user_input}, {'configurable': {'session_id': '[0]'}})

def invoke_search_chain(user_input:str):
    choose_model()
    vectorstore = get_vectordb_web()
    search = get_google_api()
    llm = st.session_state.model
    web_research_retriever = WebResearchRetriever.from_llm(
        llm=llm,
        vectorstore=vectorstore,
        search=search,
        text_splitter=RecursiveCharacterTextSplitter(chunk_size=300, chunk_overlap=30),
        num_search_results=3,
        prompt=get_search_query_prompt(),
    )
    document_prompt, question_prompt, combine_prompt = get_qa_source_prompt()
    qa_chain = RetrievalQAWithSourcesChain.from_llm(
        llm=llm,
        document_prompt = document_prompt,
        question_prompt = question_prompt,
        combine_prompt = combine_prompt,
        retriever=web_research_retriever,
        max_tokens_limit = 1000,
        verbose=global_log_verbosity,
    )
    result = qa_chain.invoke(input=user_input)
    return result['answer']

def invoke_document_chain(user_input:str):
    choose_model()
    vectorstore = get_vectordb_doc()
    llm = st.session_state.model
    retriever = vectorstore.as_retriever(search_kwargs={"k": st.session_state.doc_search_k})
    document_prompt = get_qa_prompt()
    qa_chain = RetrievalQA.from_llm(
        llm=llm,
        prompt = document_prompt,
        retriever=retriever,
        verbose=global_log_verbosity,
    )
    result = qa_chain.invoke(input=user_input)
    return result['result']

def output_messages(container):
    """
    Output session messages.
    """
    messages = st.session_state.messages
    for message in messages:
        if isinstance(message, AIMessage):
            with container.chat_message('assistant'):
                container.markdown(message.content)
        elif isinstance(message, HumanMessage):
            with container.chat_message('user'):
                container.markdown(message.content)
        else:
            pass

def do_chat_mode(container):
    chat_container = st.container(height=400, border=True)
    with container:
        if user_input := container.chat_input('Ask me anything!', key='chat'):
            with st.spinner('Chat Agent is typing ...'):
                response = invoke_chat_chain(user_input=user_input)
                human_message = HumanMessage(content=user_input)
                ai_message = AIMessage(content=response)
                st.session_state.messages.append(human_message)
                st.session_state.messages.append(ai_message)
    output_messages(chat_container)

def do_search_mode(container):
    st.session_state.use_googlesearch = container.checkbox('Google検索')
    st.session_state.search_k = st.slider('最大検索数:', min_value=1, max_value=10, value=5, step=1)

    chat_container = st.container(height=400, border=True)
    with container:
        if user_input := container.chat_input('Ask me anything!', key='search'):
            with st.spinner('Chat Agent is typing ...'):
                response = invoke_search_chain(user_input=user_input)
                human_message = HumanMessage(content=user_input)
                ai_message = AIMessage(content=response)
                st.session_state.history.add_user_message(human_message)
                st.session_state.history.add_ai_message(ai_message)
                st.session_state.messages.append(human_message)
                st.session_state.messages.append(ai_message)
    output_messages(chat_container)

def do_document_mode(container):
    upload_file = container.file_uploader(label='ドキュメント:', type=['pdf', 'txt', 'text', 'xlsx', 'csv', 'html', 'ppt', 'pptx'])
    if upload_file != st.session_state.upload_file:
        st.session_state.upload_file = upload_file
        embeddig_file(st.session_state.upload_file)
    st.session_state.doc_search_k = st.slider('最大検索数:', min_value=1, max_value=10, value=3, step=1)

    chat_container = st.container(height=400, border=True)
    with container:
        if st.session_state.upload_file:
            if user_input := container.chat_input('Ask me about the file', key='doc'):
                with st.spinner('Chat Agent is typing ...'):
                    response = invoke_document_chain(user_input=user_input)
                    human_message = HumanMessage(content=user_input)
                    ai_message = AIMessage(content=response)
                    st.session_state.history.add_user_message(human_message)
                    st.session_state.history.add_ai_message(ai_message)
                    st.session_state.messages.append(human_message)
                    st.session_state.messages.append(ai_message)
        else:
            st.write('First, upload a file')
    output_messages(chat_container)

def main():
    init_page()
    init_options()
    init_messages()
    tab_chat, tab_search, tab_doc = st.tabs(['チャット', '検索', '文書'])
    with tab_chat:
        do_chat_mode(tab_chat)

    with tab_search:
        do_search_mode(tab_search)

    with tab_doc:
        do_document_mode(tab_doc)

if __name__ == '__main__':
    main()
